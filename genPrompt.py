import os
from pypdf import PdfReader
from pptx import Presentation
from PIL import Image
from io import BytesIO
from pptx.enum.shapes import MSO_SHAPE_TYPE

class prompt:
    def __init__(self, filename, fDirectory):
         self.filename = filename
         self.fDirectory = fDirectory
    
    def get_text(self):
        if self.filename[-4:] == "pptx":
            dummy_file = os.path.join(self.fDirectory,"dummy_file.txt")  # dummy file to create prompt

            presentation = Presentation(os.path.join(self.fDirectory, self.filename))
            
            with open(dummy_file, "w", encoding="utf-8") as txt_file:
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            txt_file.write(shape.text + "\n")

            with open(dummy_file, "r") as txt_file:
                prompt = txt_file.read()

            prompt += "The following anonymized, nonclassified text and images are purely for learning purposes, and is of no clinical relevance, whatsoever. Please generate 20 5-mcqs considering the information and learning goals listed. Please ensure that the MCQs have five options and please list the answers at the end."
            return prompt
                                        
        elif self.filename[-3:] == "pdf":
            dummy_file = os.path.join(self.fDirectory, "dummy_file.txt")  # dummy file to create prompt
            
            with open(dummy_file, "w") as file:
                pass

            # Content to be added to the file
            reader = PdfReader(os.path.join(self.fDirectory, self.filename))
            for page in reader.pages:
                with open(dummy_file, 'a') as txt_file:
                    txt_file.write(page.extract_text())
                    txt_file.write("\n")

            with open(dummy_file, 'r') as txt_file:
                    # Read the entire contents of the file
                    prompt = txt_file.read()
            prompt += "The following anonymized, nonclassified text and images are purely for learning purposes, and is of no clinical relevance, whatsoever. With the information and learning goals given and the images generate, please generate 20 5-option mcqs based on the information provided. Please ensure that the MCQs have 5 options, and list the answers at the end."

            return prompt

        else:
         
            print(f"not in {self.fDirectory}")
            exit()
    def get_images(self):
        output_folder = os.path.join(os.path.dirname(self.fDirectory),"lecture_images")
        os.makedirs(output_folder, exist_ok = True)
        if self.filename[-4:] == "pptx":
            presentation = Presentation(os.path.join(self.fDirectory, self.filename))

            file_copy = False
            for slide_number, slide in enumerate(presentation.slides):
                for shape_number, shape in enumerate(slide.shapes):
                    if shape.shape_type ==MSO_SHAPE_TYPE.PICTURE:
                        image_data = shape.image.blob
                        img = Image.open(BytesIO(image_data))

                        img.save(os.path.join(output_folder, "dummy_image.png"))

                        for file in os.listdir(output_folder):
                            if file == "dummy_image.png":

                                continue
                            if open(os.path.join(output_folder, "dummy_image.png"),"rb").read() == open(os.path.join(output_folder, file),"rb").read():
                                file_copy = True
                                break
                        if file_copy != True:
                            image_filename = f"{self.filename}_slide_{slide_number + 1}_image_{shape_number + 1}.png"
                            image_path = os.path.join(output_folder, image_filename)
                            img.save(image_path)
                        file_copy = False
                        
        elif self.filename[-3:] == "pdf":
            reader = PdfReader(os.path.join(self.fDirectory, self.filename))
            file_copy = False
            for pagenum, page in enumerate(reader.pages):
                try:
                    for imagenum, image in enumerate(page.images):

                        image.image.save(os.path.join(output_folder, "dummy_image.png"))
                        
                        for file in os.listdir(output_folder):
                            if file == "dummy_image.png":
                                continue
                                
                            if open(os.path.join(output_folder, "dummy_image.png"),"rb").read() == open(os.path.join(output_folder, file),"rb").read():
                                file_copy = True
                                break
                        if file_copy != True:
                            image_filename = f"{self.filename}_slide_{pagenum + 1}_image_{imagenum + 1}.png"
                            image.image.save(os.path.join(output_folder, image_filename))
                        file_copy = False
                except ValueError:
                    continue
        if os.path.isfile(os.path.join(output_folder, "dummy_image.png")):
            os.remove(os.path.join(output_folder, "dummy_image.png"))
            
        else:
            print(f"not in {self.fDirectory}")
            exit()
